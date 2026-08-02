"""발표자료가 실제로 읽히는 상태인지 검사합니다.

이 서버에는 파워포인트도 LibreOffice도 없어 슬라이드를 눈으로 볼 수
없습니다. 그렇다고 "아마 괜찮겠지" 하고 제출할 수는 없습니다. 글자가
상자를 넘치거나 도형끼리 겹치면 심사위원 화면에서 그대로 깨져 보입니다.

그래서 만들어진 pptx를 다시 열어 세 가지를 셉니다.

  1. 슬라이드 밖으로 나간 도형
  2. 글자가 상자보다 길어 넘칠 것으로 보이는 곳
  3. 글상자끼리 겹치는 곳

글자 길이는 어림으로 잽니다. 한글은 글꼴 크기만큼, 영문·숫자는 그 절반쯤을
가로로 차지합니다. 정확하지는 않지만 **넘칠 것 같은 곳을 집어내는 데는**
충분합니다. 여유를 15% 두고 그보다 넘치면 알립니다.

실행:
    python -m scripts.check_deck
"""
from __future__ import annotations

import os
import sys
import unicodedata

from pptx import Presentation
from pptx.util import Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(ROOT, "..", "docs", "꿀비_기술설명서.pptx")

SLACK = 1.15          # 이만큼까지는 넘쳐도 눈에 안 띕니다
MIN_OVERLAP = 0.15    # 넓이의 15% 넘게 겹치면 알립니다


def _width(text: str, pt: float) -> float:
    """글자가 가로로 차지하는 폭을 어림합니다 (pt).

    한글·한자·일본어는 정사각형에 가깝고(글꼴 크기와 폭이 거의 같음),
    영문 소문자와 숫자는 절반쯤입니다. 이 둘만 갈라도 어림은 됩니다.
    """
    w = 0.0
    for ch in text:
        if unicodedata.east_asian_width(ch) in ("W", "F"):
            w += pt
        elif ch == " ":
            w += pt * 0.28
        else:
            w += pt * 0.52
    return w


def _lines(text: str, pt: float, box_pt: float) -> int:
    """줄바꿈까지 헤아려 몇 줄이 될지 셉니다."""
    if box_pt <= 0:
        return 999
    n = 0
    for para in text.split("\n"):
        n += max(1, -(-_width(para, pt) // box_pt))   # 올림
    return int(n)


def check(path: str) -> int:
    prs = Presentation(path)
    SW, SH = prs.slide_width, prs.slide_height
    problems = 0

    for i, slide in enumerate(prs.slides, 1):
        boxes = []

        for sh in slide.shapes:
            # 1. 슬라이드 밖으로 나갔는가
            if sh.left is None or sh.top is None:
                continue
            r, b = sh.left + (sh.width or 0), sh.top + (sh.height or 0)
            if sh.left < -Emu(1) or sh.top < -Emu(1) or r > SW + Emu(1) or b > SH + Emu(1):
                over_r = max(0, r - SW) / 914400
                over_b = max(0, b - SH) / 914400
                print(f"  [{i}장] 슬라이드 밖 — {sh.shape_type}, "
                      f"오른쪽 {over_r:.2f}in · 아래 {over_b:.2f}in 초과")
                problems += 1

            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue

            tf = sh.text_frame
            txt = tf.text
            pts = [r_.font.size.pt for p in tf.paragraphs for r_ in p.runs
                   if r_.font.size]
            pt = max(pts) if pts else 18.0
            spacing = next((p.line_spacing for p in tf.paragraphs if p.line_spacing), 1.25)

            box_w_pt = (sh.width or 0) / 12700
            box_h_pt = (sh.height or 0) / 12700
            need = _lines(txt, pt, box_w_pt) * pt * float(spacing)

            # 2. 글자가 상자보다 긴가
            if need > box_h_pt * SLACK:
                print(f"  [{i}장] 글자 넘침 — {need/pt:.1f}줄 필요, "
                      f"상자는 {box_h_pt/pt:.1f}줄  «{txt[:38].replace(chr(10),' ')}…»")
                problems += 1

            boxes.append((sh.left, sh.top, sh.width or 0, sh.height or 0, txt))

        # 3. 글상자끼리 겹치는가
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                ax, ay, aw, ah, at = boxes[a]
                bx, by, bw, bh, bt = boxes[b]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox <= 0 or oy <= 0:
                    continue
                area = ox * oy
                small = min(aw * ah, bw * bh)
                if small and area / small > MIN_OVERLAP:
                    print(f"  [{i}장] 글상자 겹침 {area/small*100:.0f}% — "
                          f"«{at[:22].replace(chr(10),' ')}» ↔ "
                          f"«{bt[:22].replace(chr(10),' ')}»")
                    problems += 1

    print(f"\n슬라이드 {len(prs.slides)}장 · 문제 {problems}건")
    return problems


def main() -> int:
    path = os.path.abspath(DECK)
    if not os.path.exists(path):
        print(f"파일이 없습니다 — {path}")
        return 1
    print(f"검사 — {os.path.basename(path)}")
    return 0 if check(path) == 0 else 2


if __name__ == "__main__":
    sys.exit(main())
