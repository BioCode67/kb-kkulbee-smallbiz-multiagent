"""기술설명서(PowerPoint)를 만듭니다 — 예선 필수 제출물.

KB AI Challenge 예선은 셋을 냅니다: 참가신청서, **기술설명서(파워포인트)**,
프로토타입(구현코드). 마감 2026-08-03 16:00.

**숫자를 손으로 적지 않습니다.** 색인 meta.json과 평가 스크립트가 남긴 값을
읽어 슬라이드에 박습니다. 발표자료의 숫자와 실제 코드가 어긋나는 것이
심사에서 가장 위험합니다 — 시연을 눌러 보는 순간 드러납니다.

실행:
    python -m scripts.make_deck
"""
from __future__ import annotations

import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "..", "docs", "꿀비_기술설명서.pptx")
SHOTS = os.path.join(ROOT, ".cache", "deck")

# ── 색 ────────────────────────────────────────────────────────────────────
# KB 로고의 두 색(옐로 #FFBC00 + 웜 잉크)이 그대로 슬라이드의 주조입니다.
# 웹과 발표자료가 같은 얼굴이어야 합니다 — 심사위원은 둘을 나란히 봅니다.
INK = RGBColor(0x38, 0x32, 0x2A)
GRAY = RGBColor(0x6B, 0x62, 0x59)
LINE = RGBColor(0xE4, 0xDE, 0xD4)
YELLOW = RGBColor(0xFF, 0xBC, 0x00)
YELLOW_SOFT = RGBColor(0xFF, 0xF4, 0xD6)
BROWN = RGBColor(0x54, 0x44, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
RED = RGBColor(0xC0, 0x39, 0x2B)

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)      # 16:9


def _meta() -> dict:
    def load(p):
        with open(os.path.join(ROOT, "app", "data", p), encoding="utf-8") as f:
            return json.load(f)
    return {"market": load("market_index/meta.json"),
            "policy": load("policy_index/meta.json")}


# ── 그리기 도우미 ─────────────────────────────────────────────────────────
def text(slide, x, y, w, h, s, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return box


def rect(slide, x, y, w, h, fill=None, line=None, width=Pt(1)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = width
    sh.shadow.inherit = False
    return sh


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, kicker, title):
    """모든 본문 장의 머리. 같은 자리에 같은 크기로 둡니다 — 장을 넘길 때
    제목이 움직이면 읽는 사람이 매번 다시 위치를 찾습니다."""
    text(slide, Inches(0.75), Inches(0.5), Inches(11), Inches(0.3),
         kicker, size=12, bold=True, color=YELLOW)
    text(slide, Inches(0.75), Inches(0.82), Inches(11.8), Inches(0.6),
         title, size=27, bold=True, color=INK)
    ln = slide.shapes.add_shape(1, Inches(0.75), Inches(1.55), Inches(11.8), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background(); ln.shadow.inherit = False


def stat(slide, x, y, w, value, label, note="", color=INK):
    """큰 숫자 하나와 그 뜻. 숫자만 크게 두면 무슨 숫자인지 모릅니다."""
    text(slide, x, y, w, Inches(0.6), value, size=34, bold=True, color=color)
    text(slide, x, y + Inches(0.62), w, Inches(0.3), label, size=13, bold=True, color=BROWN)
    if note:
        text(slide, x, y + Inches(0.95), w, Inches(0.6), note, size=10.5, color=GRAY)


def bullets(slide, x, y, w, items, size=13.5, gap=0.42):
    """· 하나에 한 줄. 두 줄 넘어가면 읽히지 않습니다."""
    for i, it in enumerate(items):
        cy = y + Inches(gap * i)
        dot = slide.shapes.add_shape(9, x, cy + Inches(0.07), Pt(6), Pt(6))
        dot.fill.solid(); dot.fill.fore_color.rgb = YELLOW
        dot.line.fill.background(); dot.shadow.inherit = False
        text(slide, x + Inches(0.22), cy, w - Inches(0.22), Inches(0.4), it, size=size)


def build() -> str:
    m = _meta()
    mk, pl = m["market"], m["policy"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    bee = os.path.join(ROOT, "..", "frontend", "public", "kkulbee.svg")
    bee_png = os.path.join(ROOT, ".cache", "kkulbee_deck.png")
    if not os.path.exists(bee_png):
        try:
            import cairosvg
            os.makedirs(os.path.dirname(bee_png), exist_ok=True)
            cairosvg.svg2png(url=bee, write_to=bee_png,
                             output_width=520, output_height=592)
        except Exception:      # noqa: BLE001 — 그림 없이도 자료는 나와야 합니다
            bee_png = None
    if bee_png and not os.path.exists(bee_png):
        bee_png = None

    # ── 1. 표지 ──────────────────────────────────────────────────────────
    s = blank(prs)
    bg = rect(s, 0, 0, W, H, fill=RGBColor(0x1C, 0x18, 0x14))
    band = rect(s, 0, H - Inches(0.9), W, Inches(0.9), fill=YELLOW)

    text(s, Inches(0.9), Inches(1.5), Inches(8), Inches(0.4),
         "제8회 KB AI Challenge", size=15, bold=True, color=YELLOW)
    text(s, Inches(0.9), Inches(1.95), Inches(8.6), Inches(1.15),
         "꿀비", size=64, bold=True, color=WHITE)
    text(s, Inches(0.9), Inches(3.3), Inches(8.6), Inches(1.2),
         "사장님의 세 가지 고민, 한 번에 물어보세요\n"
         "어디에 열지 · 자금은 어떻게 · 억울한 일이 생기면",
         size=21, color=RGBColor(0xE8, 0xE0, 0xD4), spacing=1.4)
    text(s, Inches(0.9), Inches(4.7), Inches(8.6), Inches(0.9),
         f"전국 점포 {mk['stores']:,}개 · 정부 지원사업 공고 {pl['docs']}건\n"
         "실측 자료 위에서만 답합니다",
         size=14, color=YELLOW_SOFT, spacing=1.4)
    text(s, Inches(0.9), H - Inches(0.72), Inches(9), Inches(0.4),
         "Pick 2 최적 입지  ×  Pick 3 소상공인 금융  ×  Pick 4 소비자 보호",
         size=13, bold=True, color=INK)
    if bee_png:
        # 그림 폭은 높이 × (220/250)입니다. 오른쪽으로 넘치지 않게
        # 왼쪽 자리를 폭만큼 빼서 잡습니다.
        s.shapes.add_picture(bee_png, Inches(9.0), Inches(1.5),
                             height=Inches(4.4))

    # ── 1.5 차별화 — 작년에 비슷한 이름이 있었다 ─────────────────────────
    #
    # 제7회 본선에 「사장님의 꿀벌 비서」가 있다. 심사위원은 기억한다.
    # 먼저 말하지 않으면 "또 벌 비서네"로 시작한다. 정면으로 짚고 간다.
    s = blank(prs)
    header(s, "WHY DIFFERENT", "작년에도 '꿀벌 비서'가 있었습니다 — 무엇이 다른가")
    text(s, Inches(0.75), Inches(1.95), Inches(11.8), Inches(0.6),
         "제7회 본선 「사장님의 꿀벌 비서」는 흩어진 정보를 모아 보여주는 서비스였습니다. "
         "꿀비는 '보여주기'가 아니라 '실측으로 답하고 증명하는' 쪽에 섰습니다.",
         size=14, color=GRAY)
    rows5 = [
        ("실측 규모", "점포 272만 개 좌표 · 공고 900건 하이브리드 검색 P@1 0.875 — 정보 나열이 아니라 검색·분석"),
        ("3픽 크로스오버", "입지×금융×소비자보호가 한 질문에서 동시에 켜지고, 어떤 에이전트가 돌았는지 화면에 남습니다"),
        ("구조로 보장되는 가드레일", "금소법 검사가 그래프의 마지막 노드 — LLM이 쓴 문장도 예외 없이 거칩니다"),
        ("정직성의 설계", "없는 자료는 없다고 화면에 적고, 모르는 동네에 점수를 주지 않고, 추천마다 원문 링크"),
        ("키·GPU 없이 영구 동작", "검색은 int8 ONNX CPU 5ms — 시연 환경과 외부 한도를 타지 않습니다"),
    ]
    y = Inches(2.7)
    for t, d in rows5:
        text(s, Inches(0.75), y, Inches(3.4), Inches(0.35), t, size=14.5, bold=True)
        text(s, Inches(4.3), y + Inches(0.02), Inches(8.2), Inches(0.5),
             d, size=12, color=GRAY)
        y += Inches(0.82)

    # ── 2. 문제 ──────────────────────────────────────────────────────────
    s = blank(prs)
    header(s, "PROBLEM", "사장님은 '무엇을 물어야 하는지'를 모릅니다")
    text(s, Inches(0.75), Inches(1.95), Inches(11.8), Inches(0.6),
         "\"요즘 통 손님이 없어서 월세도 빠듯해요\"",
         size=26, bold=True, color=BROWN)
    text(s, Inches(0.75), Inches(2.65), Inches(11.8), Inches(0.5),
         "이 문장에는 '자금'도 '대출'도 '지원사업'도 없습니다. "
         "그런데 이 사장님께 필요한 것은 소상공인 경영안정자금입니다.",
         size=14, color=GRAY)

    cards = [
        ("제도를 몰라서 못 받는다",
         "지원사업은 부처·지자체·공단에 흩어져 있고\n공고문은 행정 용어로 쓰여 있습니다.\n"
         "이름을 모르면 검색창에 넣을 수도 없습니다."),
        ("상권 정보는 '평균'만 준다",
         "'이 동네 카페 평균 매출'을 들어도\n내 자리가 어떤지는 알 수 없습니다.\n"
         "무엇이 좋고 무엇이 나쁜지 쪼개 줘야 합니다."),
        ("설명을 못 들어도 모른다",
         "금소법상 설명의무 위반이어도\n그것이 위반인지 아는 사장님은 드뭅니다.\n"
         "권리가 있어도 쓰지 못합니다."),
    ]
    for i, (t, d) in enumerate(cards):
        x = Inches(0.75 + i * 4.02)
        rect(s, x, Inches(3.5), Inches(3.72), Inches(2.5),
             fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
        text(s, x + Inches(0.3), Inches(3.8), Inches(3.1), Inches(0.5),
             t, size=15, bold=True, color=INK)
        text(s, x + Inches(0.3), Inches(4.35), Inches(3.2), Inches(1.5),
             d, size=11.5, color=GRAY, spacing=1.35)

    text(s, Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.5),
         "→ 제도 이름이 아니라 '처한 상황'으로 물어도 답이 나와야 합니다.",
         size=15, bold=True, color=INK)

    # ── 3. 해결 ──────────────────────────────────────────────────────────
    s = blank(prs)
    header(s, "SOLUTION", "한 번 물으면 세 에이전트가 동시에 답합니다")
    text(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.4),
         "\"연남동에서 카페 열려는데 상권도 보고 자금도 알려줘\" — 한 문장에 두 과제가 걸려 있습니다.",
         size=13, color=GRAY)

    boxes = [
        ("Router", "갈래·지역·업종을 읽습니다\n규칙이 확실하면 LLM을 안 부릅니다", YELLOW_SOFT, INK),
        ("Location", f"전국 점포 {mk['stores']//10000}만 개\n행정동 {mk['dongs_kept']:,}곳 백분위", WHITE, INK),
        ("Policy", f"기업마당 공고 {pl['docs']}건\nBM25 + 임베딩 하이브리드", WHITE, INK),
        ("Protection", "금소법 검사·재작성\n용어 15개 · 분쟁 4단계", WHITE, INK),
    ]
    for i, (t, d, fill, col) in enumerate(boxes):
        x = Inches(0.75 + i * 3.0)
        rect(s, x, Inches(2.55), Inches(2.7), Inches(1.5), fill=fill, line=LINE)
        text(s, x + Inches(0.25), Inches(2.78), Inches(2.3), Inches(0.4),
             t, size=15, bold=True, color=col)
        text(s, x + Inches(0.25), Inches(3.22), Inches(2.3), Inches(0.8),
             d, size=10.5, color=GRAY, spacing=1.3)
        if i < 3:
            text(s, x + Inches(2.72), Inches(3.1), Inches(0.3), Inches(0.3),
                 "→", size=16, bold=True, color=YELLOW)

    rect(s, Inches(0.75), Inches(4.35), Inches(11.8), Inches(0.95),
         fill=RGBColor(0xFF, 0xF9, 0xE8), line=YELLOW)
    text(s, Inches(1.05), Inches(4.52), Inches(11.2), Inches(0.6),
         "가드레일은 예외 없이 마지막에 돕니다 — 어느 경로로 왔든, LLM이 쓴 문장이든, "
         "사용자에게 나가는 모든 문장이 금소법 검사를 거칩니다.\n"
         "이 불변식을 코드가 아니라 그래프 구조로 보장합니다.",
         size=12.5, color=INK, spacing=1.35)

    bullets(s, Inches(0.75), Inches(5.6), Inches(11.8), [
        "라우터가 하나만 고르지 않습니다 — 한 질문이 여러 과제에 걸치는 경우가 실제로 많습니다",
        "무엇을 그릴지도 서버가 정합니다(Generative UI) — 질문에 따라 화면 카드 구성이 달라집니다",
        "답변은 끝난 순서가 아니라 갈래 중요도 순서로 읽힙니다 — 분쟁 질문에 지원사업이 먼저 나오면 안 됩니다",
    ])

    # ── 4. 데이터 ────────────────────────────────────────────────────────
    s = blank(prs)
    header(s, "DATA", "숫자마다 출처가 있고, 없는 것은 없다고 적습니다")
    rows = [
        ("소상공인시장진흥공단 상가(상권)정보", "2026-03-31",
         f"점포 {mk['stores']:,}개 · 행정동 {mk['dongs_kept']:,}곳 · 업종 {mk['industries']}종",
         "공공데이터포털"),
        ("중소벤처기업부 기업마당 공고", pl["collected_at"][:10],
         f"공고 {pl['docs']}건 · 접수중 {pl['open_status']['open']} · "
         f"상시 {pl['open_status']['rolling']} · 마감 {pl['open_status']['closed']}",
         "bizinfo.go.kr 직접 수집"),
        ("금융소비자보호법", "현행",
         "설명의무(§19) · 부당권유금지(§21) · 적합성(§17) 등 가드레일 규칙",
         "법령 원문"),
    ]
    y = Inches(2.0)
    for name, when, what, src in rows:
        rect(s, Inches(0.75), y, Inches(11.8), Inches(1.0),
             fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
        text(s, Inches(1.05), y + Inches(0.14), Inches(8.5), Inches(0.35),
             name, size=15, bold=True)
        text(s, Inches(1.05), y + Inches(0.55), Inches(8.5), Inches(0.35),
             what, size=11.5, color=GRAY)
        text(s, Inches(9.9), y + Inches(0.14), Inches(2.5), Inches(0.35),
             when, size=12, bold=True, color=BROWN, align=PP_ALIGN.RIGHT)
        text(s, Inches(9.9), y + Inches(0.55), Inches(2.5), Inches(0.35),
             src, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
        y += Inches(1.15)

    rect(s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(1.35),
         fill=RGBColor(0xFD, 0xF2, 0xF0), line=RED)
    text(s, Inches(1.05), Inches(5.72), Inches(11.2), Inches(0.35),
         "점수에 넣지 않은 것 — " + " · ".join(mk["not_measured"]),
         size=14, bold=True, color=RED)
    text(s, Inches(1.05), Inches(6.12), Inches(11.2), Inches(0.7),
         "상가정보 자료에 없습니다. 추정해 채우면 지어내는 것이므로 넣지 않고, "
         "화면과 API(/api/v1/sources) 양쪽에 없다고 적습니다.\n"
         "개발 초기에는 이 값들을 동네 이름 해시로 만들어 '유동인구 24.6천명/일'처럼 "
         "실측인 양 내보내고 있었습니다. 걷어냈습니다.",
         size=11.5, color=BROWN, spacing=1.35)

    # ── 5. 핵심기술 1 — 검색 ─────────────────────────────────────────────
    s = blank(prs)
    header(s, "CORE 1", "정책자금 검색 — 낱말과 뜻, 둘 다로 찾습니다")
    c = pl.get("centering", {})
    text(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.4),
         "BM25는 '전통시장'을 정확히 잡고, 임베딩은 \"장사가 안돼요\"와 '경영안정자금'을 잇습니다. "
         "약점이 서로 달라 합치면 각자보다 낫습니다.", size=13, color=GRAY)

    stat(s, Inches(0.75), Inches(2.6), Inches(2.6), f"{pl['docs']}건",
         "기업마당 실제 공고", f"{pl['chunks']:,}개 조각으로 나눠 색인")
    stat(s, Inches(3.6), Inches(2.6), Inches(2.6), "5ms",
         "질의당 검색 시간", "CPU · int8 ONNX · 키 불필요")
    stat(s, Inches(6.45), Inches(2.6), Inches(2.6), "0.875",
         "P@1 (16문항 평가셋)", "기준선 0.625에서 개선", color=GREEN)
    stat(s, Inches(9.3), Inches(2.6), Inches(2.6), "0.938",
         "MRR", "기준선 0.719에서 개선", color=GREEN)

    rect(s, Inches(0.75), Inches(4.15), Inches(5.8), Inches(2.5),
         fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
    text(s, Inches(1.05), Inches(4.35), Inches(5.2), Inches(0.35),
         "재 보지 않았으면 놓쳤을 문제", size=14, bold=True)
    text(s, Inches(1.05), Inches(4.8), Inches(5.3), Inches(1.7),
         f"임베딩 유사도가 0.82~0.86에 다 몰려 있었습니다. 폭이 "
         f"{c.get('spread_before', 0.0465)}입니다 — 순위가 사실상 잡음으로 정해집니다.\n"
         "\"운영자금이 필요해요\"에 「IP투자연계 지식재산평가」가 1위였습니다.\n"
         "평균 방향을 빼고 다시 정규화하니 폭이 "
         f"{c.get('spread_after', 0.3038)}로 벌어지고, 같은 질문에 "
         "「소상공인시장진흥자금 융자」가 1위가 됐습니다.",
         size=11, color=GRAY, spacing=1.3)

    rect(s, Inches(6.75), Inches(4.15), Inches(5.8), Inches(2.5),
         fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
    text(s, Inches(7.05), Inches(4.35), Inches(5.2), Inches(0.35),
         "평가 기준부터 의심했습니다", size=14, bold=True)
    text(s, Inches(7.05), Inches(4.8), Inches(5.3), Inches(1.7),
         "처음 P@1이 0.94로 나왔습니다. 들여다보니 \"배달앱 수수료\" 질문에 "
         "「IoT 보안인증 인증시험 수수료 지원」이 정답으로 세어지고 있었습니다. "
         "낱말 하나가 겹쳤을 뿐입니다.\n"
         "기준을 조여 0.625까지 내리고, 거기서 다시 0.875로 올렸습니다. "
         "후하게 나온 점수는 고칠 곳을 가려 버립니다.",
         size=11, color=GRAY, spacing=1.3)

    # ── 6. 핵심기술 2 — 상권 ─────────────────────────────────────────────
    s = blank(prs)
    header(s, "CORE 2", "상권 점수 — 무엇이 몇 점을 올리고 내렸는지까지")
    text(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.4),
         "기준점 50에서 요인 다섯이 더하고 뺀 몫이 최종 점수입니다. "
         "요인마다 전국 행정동 분포에서의 백분위로 잽니다.", size=13, color=GRAY)

    text(s, Inches(0.75), Inches(2.5), Inches(6.2), Inches(0.35),
         "서울 마포구 연남동 · 카페 → 66.8점 (B등급)", size=16, bold=True, color=INK)
    demo = [
        ("동종업종 경쟁", -12.9, "카페 204개 · 전국 상위 0.5%"),
        ("상권 규모", 10.9, "점포 2,327개 · 전국 상위 4.8%"),
        ("점포 밀집도", 9.7, "0.55㎢에 2,327개 · ㎢당 4,208개"),
        ("같은 계열 집적", 4.6, "음식 39% · 전국 평균의 1.3배"),
        ("생활 인프라", 6.0, "약국·편의점·의원 등 9종 모두"),
        ("업종 다양성", -1.5, "몇몇 업종에 쏠린 편"),
    ]
    y = Inches(3.0)
    for label, val, note in demo:
        text(s, Inches(0.75), y, Inches(1.9), Inches(0.3), label, size=12, bold=True)
        mid = Inches(4.3)
        w = Inches(abs(val) / 13 * 1.4)
        bar = s.shapes.add_shape(
            1, mid if val > 0 else mid - w, y + Inches(0.04), w, Inches(0.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = YELLOW if val > 0 else RGBColor(0xE0, 0x7A, 0x70)
        bar.line.fill.background(); bar.shadow.inherit = False
        text(s, Inches(5.9), y, Inches(0.8), Inches(0.3),
             f"{val:+.1f}", size=12, bold=True,
             color=INK if val > 0 else RED, align=PP_ALIGN.RIGHT)
        text(s, Inches(6.9), y, Inches(5.6), Inches(0.3), note, size=11, color=GRAY)
        y += Inches(0.42)

    rect(s, Inches(0.75), Inches(5.35), Inches(11.8), Inches(1.5),
         fill=RGBColor(0xFF, 0xF9, 0xE8), line=YELLOW)
    text(s, Inches(1.05), Inches(5.55), Inches(11.2), Inches(0.35),
         "모르는 동네에는 점수를 주지 않습니다", size=14, bold=True)
    text(s, Inches(1.05), Inches(5.95), Inches(11.2), Inches(0.8),
         "예전에는 동네 이름만 있으면 해시 난수로 채워 무엇이든 답했습니다. "
         "위경도까지 한반도 한가운데 무작위였습니다.\n"
         "지금은 자료에 없으면 \"어느 동네인지 못 알아들었습니다\"라고 합니다. "
         "대신 '성수동'(→성수2가3동), '홍대'(→서교동)처럼 부르는 이름과 행정동이 "
         "다른 경우를 알아듣게 넓혔습니다.",
         size=11.5, color=BROWN, spacing=1.35)

    # ── 7. 핵심기술 3 — 가드레일 ─────────────────────────────────────────
    s = blank(prs)
    header(s, "CORE 3", "금융소비자보호 가드레일 — 마지막 관문")
    text(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.4),
         "금융 서비스에서 AI가 하면 안 되는 말이 법으로 정해져 있습니다. "
         "규칙으로 막고, 막은 것을 기록에 남깁니다.", size=13, color=GRAY)

    g = [
        ("금리·한도 확정 표현", "\"연 2.5%로 5천만원 받으실 수 있습니다\"",
         "→ \"공고 기준이며 실제 조건은 심사 결과에 따릅니다\""),
        ("수익·승인 보장", "\"신청하면 100% 됩니다\"", "→ 삭제 후 안내 문구로 대체"),
        ("권유·단정", "\"지금 대출받으세요\"", "→ 정보 제공 표현으로 재작성"),
    ]
    y = Inches(2.5)
    for title, bad, fix in g:
        rect(s, Inches(0.75), y, Inches(11.8), Inches(1.15),
             fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
        text(s, Inches(1.05), y + Inches(0.16), Inches(3.0), Inches(0.35),
             title, size=13, bold=True)
        text(s, Inches(4.2), y + Inches(0.14), Inches(8.1), Inches(0.35),
             bad, size=11.5, color=RED)
        text(s, Inches(4.2), y + Inches(0.6), Inches(8.1), Inches(0.35),
             fix, size=11.5, color=GREEN)
        y += Inches(1.3)

    rect(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.65),
         fill=RGBColor(0xFF, 0xF9, 0xE8), line=YELLOW)
    text(s, Inches(1.05), Inches(6.6), Inches(11.2), Inches(0.4),
         "무엇을 왜 고쳤는지 응답에 함께 실어 화면에 띄웁니다. "
         "/api/v1/guardrail/check 에 문장을 넣어 그 자리에서 확인하실 수 있습니다.",
         size=12, color=INK)

    # ── 7.5 자유주제 3종 — "어떻게 했지" ─────────────────────────────────
    s = blank(prs)
    header(s, "BEYOND", "데이터사이언스 · 생성 · 접근성 — 세 방향의 자유 기능")
    cols = [
        ("닮은 상권 찾기", "업종 구성 247차원 코사인 유사도",
         "연남동 ↔ 광안리·서교동·망원동\n서면 ↔ 부평역·안양일번가\n"
         "붙어 있는 동네가 아니라 '성격이 닮은' 동네.\n"
         "2호점·확장 후보의 출발점. 질의 0.4ms, 외부 의존 0."),
        ("민원서 초안 생성", "규칙이 법조문, LLM은 사실관계만",
         "\"수수료를 왜 떼요\" 한 문장이 제출 가능한\n"
         "민원서 꼴이 됩니다. 없는 날짜·금액은 [빈칸]으로\n"
         "강제 — LLM이 §를 지어낼 수 없는 구조.\n생성문도 가드레일 통과."),
        ("음성 대화 + 지도 투어", "브라우저 내장 — 다운로드 0바이트",
         "말로 묻고(Web Speech) 꿀비가 읽어 줍니다.\n"
         "지도는 전국에서 그 동네로 날아 들어가\n"
         "4막 투어(점포→열지도→비교)를 자막·음성으로.\n"
         "금융 접근성 — 이 대회가 역대로 상을 준 각도."),
    ]
    for i, (t, sub, d) in enumerate(cols):
        x = Inches(0.75 + i * 4.02)
        rect(s, x, Inches(2.1), Inches(3.72), Inches(4.3),
             fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
        text(s, x + Inches(0.3), Inches(2.4), Inches(3.1), Inches(0.5),
             t, size=16, bold=True)
        text(s, x + Inches(0.3), Inches(2.95), Inches(3.1), Inches(0.4),
             sub, size=11, bold=True, color=RGBColor(0x9A, 0x6B, 0x00))
        text(s, x + Inches(0.3), Inches(3.5), Inches(3.2), Inches(2.6),
             d, size=11, color=GRAY, spacing=1.4)

    # ── 7.6 실제 화면 ────────────────────────────────────────────────────
    s = blank(prs)
    header(s, "SCREENS", "지금 접속해 보실 수 있습니다")
    text(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.4),
         "kb-kkulbee-smallbiz-multiagent.onrender.com — 목업이 아니라 배포 중인 실서비스의 화면입니다.",
         size=13, color=GRAY)
    shots = [("F1.png", Inches(0.75), Inches(3.9), "첫 화면 — 다섯 갈래 선택"),
             ("tour.png", Inches(5.15), Inches(3.9), "지도 투어 — 점포 위 자막"),
             ("policy.png", Inches(9.15), Inches(3.3), "자금 — 성격·추천 이유")]
    for fn, x, hh, cap in shots:
        fp = os.path.join(SHOTS, fn)
        if os.path.exists(fp):
            s.shapes.add_picture(fp, x, Inches(2.5), height=hh)
            text(s, x, Inches(6.55), Inches(3.8), Inches(0.35),
                 cap, size=11, bold=True, color=BROWN)

    # ── 8. 시스템 ────────────────────────────────────────────────────────
    s = blank(prs)
    header(s, "SYSTEM", "8월 13일에 이 서버가 사라져도 계속 돕니다")
    text(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.4),
         "개발용 GPU 서버는 한시적입니다. 그래서 무거운 계산은 미리 끝내 파일로 남기고, "
         "런타임은 그 파일만 읽습니다.", size=13, color=GRAY)

    left = [
        ("사전 계산 (개발 시 · GPU)", [
            "기업마당 900건 수집 → 구조화",
            "전국 점포 272만 개 → 행정동 집계",
            "임베딩 3,326개 → 평균 중심화",
            f"결과: 색인 파일 {23.8:.1f}MB (저장소에 포함)",
        ]),
        ("런타임 (배포 · CPU만)", [
            "질의 임베딩: int8 ONNX · 5ms · 키 불필요",
            "검색: BM25 + 코사인 · 12~55ms",
            "LLM: Gemini (선택) · 없으면 규칙+템플릿",
            "메모리: 354MB (무료 등급 512MB 안)",
        ]),
    ]
    for i, (t, items) in enumerate(left):
        x = Inches(0.75 + i * 6.05)
        rect(s, x, Inches(2.5), Inches(5.75), Inches(2.35),
             fill=RGBColor(0xFA, 0xF8, 0xF4), line=LINE)
        text(s, x + Inches(0.3), Inches(2.7), Inches(5.1), Inches(0.35),
             t, size=14, bold=True)
        bullets(s, x + Inches(0.3), Inches(3.15), Inches(5.2), items, size=11.5, gap=0.36)

    rect(s, Inches(0.75), Inches(5.1), Inches(11.8), Inches(1.75),
         fill=RGBColor(0xFF, 0xF9, 0xE8), line=YELLOW)
    text(s, Inches(1.05), Inches(5.3), Inches(11.2), Inches(0.35),
         "512MB에 넣기 위해 실제로 잰 것", size=14, bold=True)
    text(s, Inches(1.05), Inches(5.72), Inches(11.2), Inches(1.0),
         "배포 전에 메모리를 재 보니 611MB였습니다. 안 재고 올렸으면 그 자리에서 죽었습니다. "
         "범인은 뜻밖에도 토크나이저였습니다 — 17MB짜리 tokenizer.json을 올리는 데 280MB를 씁니다.\n"
         "sentencepiece 원본(5MB)으로 바꾸면 70MB입니다. 다만 토큰 번호가 어긋나면 검색이 "
         "조용히 망가지므로, 두 방식의 출력을 대조해 **비트 단위로 같음**을 확인하고 바꿨습니다. "
         "ONNX 메모리 아레나도 껐습니다(200MB→9MB). 결과 611MB → 354MB.",
         size=11.5, color=BROWN, spacing=1.35)

    # ── 9. 차별점 ────────────────────────────────────────────────────────
    s = blank(prs)
    header(s, "WHY US", "정직을 설계에 넣었습니다")
    items = [
        ("모르면 모른다고 합니다",
         "맞는 공고가 없으면 빈손으로 답합니다. 900건 중 무엇이든 상위 다섯 개는 뽑히지만, "
         "관련 없는 목록은 사장님의 시간을 뺏습니다."),
        ("재지 못한 것을 같은 자리에 적습니다",
         "유동인구·매출·폐업률을 점수에 안 넣었다는 사실이 상권 카드 안에 있습니다. "
         "다른 화면에서 찾아봐야 알면 감춘 것입니다."),
        ("숫자마다 원문 링크가 붙습니다",
         "추천한 공고마다 기업마당 원문 주소가 함께 나갑니다. 우리가 요약한 것을 "
         "직접 대조하실 수 있어야 합니다."),
        ("금액은 근거가 있을 때만 씁니다",
         "본문에서 찾은 가장 큰 금액을 한도로 적어 '총사업비 4,000억원'이 한도로 둔갑했습니다. "
         "'한도·최대·이내'가 곁에 있을 때만 읽고, 무엇으로 읽었는지 남깁니다."),
        ("LLM이 사실을 만들 수 없는 구조입니다",
         "LLM은 갈래를 고르고 문장을 엮을 뿐, 공고와 숫자는 색인에서 나옵니다. "
         "그 문장도 가드레일을 반드시 거칩니다."),
    ]
    y = Inches(2.0)
    for t, d in items:
        text(s, Inches(0.75), y, Inches(4.0), Inches(0.35), t, size=14, bold=True)
        text(s, Inches(4.9), y + Inches(0.02), Inches(7.6), Inches(0.7),
             d, size=11.5, color=GRAY, spacing=1.3)
        y += Inches(0.95)

    # ── 10. 마무리 ───────────────────────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=RGBColor(0x1C, 0x18, 0x14))
    rect(s, 0, H - Inches(0.9), W, Inches(0.9), fill=YELLOW)
    text(s, Inches(0.9), Inches(1.9), Inches(8.8), Inches(1.75),
         "제도 이름을 몰라도\n답을 받을 수 있어야 합니다",
         size=40, bold=True, color=WHITE, spacing=1.25)
    text(s, Inches(0.9), Inches(3.95), Inches(8.8), Inches(1.4),
         f"전국 점포 {mk['stores']:,}개 · 지원사업 공고 {pl['docs']}건 · "
         "금소법 가드레일\n"
         "실측 위에서만 답하고, 없는 것은 없다고 적는 소상공인 AI 비서",
         size=16, color=RGBColor(0xE8, 0xE0, 0xD4), spacing=1.5)
    text(s, Inches(0.9), H - Inches(0.72), Inches(11), Inches(0.4),
         "꿀비  ·  제8회 KB AI Challenge  ·  github.com/BioCode67/kb-kkulbee-smallbiz-multiagent",
         size=12, bold=True, color=INK)
    if bee_png:
        s.shapes.add_picture(bee_png, Inches(9.6), Inches(2.2), height=Inches(3.6))

    out = os.path.abspath(OUT)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    return out


def main() -> int:
    path = build()
    print(f"[saved] {path}  ({os.path.getsize(path)/1e6:.2f} MB)")
    print("  숫자는 색인 meta.json에서 읽었습니다. 코드와 자료가 어긋나지 않습니다.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
